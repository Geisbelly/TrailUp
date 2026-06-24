from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

_DEFAULT_THEME = {
    "fundo": "#171624",
    "titulo": "#F3ECDA",
    "texto": "#E3DBC8",
    "destaque": "#C99F58",
}
_ASSETS_DIR = Path(__file__).resolve().parents[1] / "assets"


def _hex_to_rgb(value: str, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    text = str(value or "").strip().lstrip("#")
    if len(text) != 6:
        return fallback
    try:
        return tuple(int(text[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
    except Exception:
        return fallback


def _resolve_asset_path(name: str | None) -> Path | None:
    candidate = str(name or "").strip()
    if not candidate:
        return None

    raw = Path(candidate)
    if raw.is_file():
        return raw

    folders = [
        _ASSETS_DIR / "ImagensReferencia",
        _ASSETS_DIR / "images",
        _ASSETS_DIR / "imgPerfil",
    ]
    for folder in folders:
        file_path = folder / candidate
        if file_path.is_file():
            return file_path
    return None


def _merge_theme(global_theme: dict[str, Any], local_theme: dict[str, Any] | None) -> dict[str, Any]:
    merged = dict(global_theme)
    local = local_theme if isinstance(local_theme, dict) else {}
    cores = local.get("cores") if isinstance(local.get("cores"), dict) else {}
    merged.update({key: value for key, value in local.items() if key != "cores" and value})
    for key in ("fundo", "titulo", "texto", "destaque"):
        if cores.get(key):
            merged[key] = cores[key]
    return merged


def _apply_slide_background(slide: Any, theme: dict[str, Any]) -> None:
    fill = slide.background.fill
    fill.solid()
    rgb = _hex_to_rgb(str(theme.get("fundo") or ""), (23, 22, 36))
    fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def _paint_title(shape: Any, text: str, *, color_hex: str, size_pt: int) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(size_pt)
    r, g, b = _hex_to_rgb(color_hex, (243, 236, 218))
    run.font.color.rgb = RGBColor(r, g, b)


def _paint_paragraph(paragraph: Any, text: str, *, color_hex: str, size_pt: int, level: int = 0) -> None:
    paragraph.text = text
    paragraph.level = level
    r, g, b = _hex_to_rgb(color_hex, (227, 219, 200))
    if paragraph.runs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)
            run.font.color.rgb = RGBColor(r, g, b)


def _add_side_image(slide: Any, image_name: str | None) -> None:
    path = _resolve_asset_path(image_name)
    if path is None:
        return
    try:
        slide.shapes.add_picture(str(path), left=Pt(450), top=Pt(100), width=Pt(220), height=Pt(220))
    except Exception:
        # Mantém o render resiliente caso a imagem não seja compatível.
        return


def gerar_pptx(
    *,
    titulo: str,
    abertura: str,
    slides: list[dict[str, Any]],
    tema_visual: dict[str, Any] | None = None,
) -> bytes:
    prs = Presentation()
    global_theme = _merge_theme(_DEFAULT_THEME, tema_visual if isinstance(tema_visual, dict) else None)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    _apply_slide_background(slide, global_theme)
    _paint_title(slide.shapes.title, titulo or "Apresentação", color_hex=str(global_theme["titulo"]), size_pt=36)

    subtitle = slide.placeholders[1]
    _paint_paragraph(subtitle.text_frame.paragraphs[0], abertura or "", color_hex=str(global_theme["texto"]), size_pt=18)
    _add_side_image(slide, str(global_theme.get("imagem_referencia") or ""))

    for item in slides:
        if not isinstance(item, dict):
            continue

        slide_theme = _merge_theme(global_theme, item.get("tema_visual") if isinstance(item.get("tema_visual"), dict) else None)
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        _apply_slide_background(slide, slide_theme)

        title = slide.shapes.title
        _paint_title(
            title,
            str(item.get("titulo") or "Tópico"),
            color_hex=str(slide_theme["titulo"]),
            size_pt=30,
        )

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()

        subtitulo = str(item.get("subtitulo") or "").strip()
        if subtitulo:
            p_sub = tf.paragraphs[0]
            _paint_paragraph(p_sub, subtitulo, color_hex=str(slide_theme["destaque"]), size_pt=18)

        pontos = item.get("pontos") if isinstance(item.get("pontos"), list) else []
        for idx, ponto in enumerate(pontos):
            text = str(ponto).strip()
            if not text:
                continue
            if idx == 0 and not subtitulo:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            _paint_paragraph(p, text, color_hex=str(slide_theme["texto"]), size_pt=16, level=0)

        image_ref = str(item.get("imagem_referencia") or slide_theme.get("imagem_referencia") or "")
        _add_side_image(slide, image_ref)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
