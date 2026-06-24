"""Extrator de PPTX: python-pptx com fallback XML."""

from __future__ import annotations

import io
import re
import zipfile
import xml.etree.ElementTree as ET
from typing import Any

from app.ingestion.models import Block, BlockKind, FileFamily


def extract(raw: bytes, *, filename: str = "apresentacao.pptx") -> dict:
    """
    Extrai blocos estruturados de um PPTX.

    Retorna dict com:
      - family, title, blocks (list[Block]), metadata
    """
    base_title = _normalize_title(filename)

    parsed = _extract_with_python_pptx(raw=raw, filename=filename, base_title=base_title)
    if parsed is not None:
        return parsed

    return _extract_with_xml_fallback(raw=raw, filename=filename, base_title=base_title)


def _extract_with_python_pptx(*, raw: bytes, filename: str, base_title: str) -> dict | None:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except Exception:
        return None

    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception:
        return None

    blocks: list[Block] = []
    title = base_title
    block_counter = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_title_text: str | None = None
        ordered_shapes = sorted(
            slide.shapes,
            key=lambda shape: (
                int(getattr(shape, "top", 0) or 0),
                int(getattr(shape, "left", 0) or 0),
            ),
        )

        for shape in ordered_shapes:
            lines = _shape_text_lines(shape)
            if not lines:
                continue

            is_title = False
            try:
                if getattr(shape, "is_placeholder", False):
                    ph_type = shape.placeholder_format.type
                    is_title = ph_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
            except Exception:
                is_title = False

            if is_title:
                heading = lines[0]
                if slide_num == 1 and heading:
                    title = heading
                slide_title_text = heading
                block_counter += 1
                blocks.append(
                    Block(
                        block_id=f"s{slide_num}_title",
                        kind=BlockKind.SLIDE_TITLE,
                        text=heading,
                        order=block_counter,
                        source_ref={"slide": slide_num},
                        metadata={"is_title": True},
                    )
                )
                lines = lines[1:]

            for line in lines:
                if not line:
                    continue
                block_counter += 1
                blocks.append(
                    Block(
                        block_id=f"s{slide_num}_b{block_counter}",
                        kind=BlockKind.LIST_ITEM,
                        text=line,
                        order=block_counter,
                        source_ref={"slide": slide_num, "slide_title": slide_title_text},
                        metadata={},
                    )
                )

        notes_text = _extract_notes_text(slide)
        if notes_text:
            block_counter += 1
            blocks.append(
                Block(
                    block_id=f"s{slide_num}_notes",
                    kind=BlockKind.SPEAKER_NOTE,
                    text=notes_text,
                    order=block_counter,
                    source_ref={"slide": slide_num},
                    metadata={"is_note": True},
                )
            )

    return {
        "family": FileFamily.PRESENTATION,
        "title": title,
        "blocks": blocks,
        "metadata": {
            "slides": len(prs.slides),
            "filename": filename,
            "extraction_mode": "python-pptx",
        },
    }


def _extract_with_xml_fallback(*, raw: bytes, filename: str, base_title: str) -> dict:
    blocks: list[Block] = []
    title = base_title
    slide_count = 0
    block_counter = 0

    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            slide_files = _sorted_pptx_xml_names(
                n
                for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            note_files = _sorted_pptx_xml_names(
                n
                for n in zf.namelist()
                if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")
            )
            slide_count = len(slide_files)

            for slide_num, xml_name in enumerate(slide_files, start=1):
                texts = _read_xml_texts(zf=zf, xml_name=xml_name)
                if not texts:
                    continue
                for idx, text in enumerate(texts):
                    if not text:
                        continue
                    kind = BlockKind.SLIDE_TITLE if idx == 0 else BlockKind.LIST_ITEM
                    if slide_num == 1 and idx == 0:
                        title = text
                    block_counter += 1
                    blocks.append(
                        Block(
                            block_id=f"s{slide_num}_b{block_counter}",
                            kind=kind,
                            text=text,
                            order=block_counter,
                            source_ref={"slide": slide_num},
                            metadata={},
                        )
                    )

            for slide_num, xml_name in enumerate(note_files, start=1):
                notes = _read_xml_texts(zf=zf, xml_name=xml_name)
                notes_text = _join_lines(notes)
                if not notes_text:
                    continue
                block_counter += 1
                blocks.append(
                    Block(
                        block_id=f"s{slide_num}_notes_xml",
                        kind=BlockKind.SPEAKER_NOTE,
                        text=notes_text,
                        order=block_counter,
                        source_ref={"slide": slide_num},
                        metadata={"is_note": True},
                    )
                )
    except Exception:
        return {
            "family": FileFamily.PRESENTATION,
            "title": title,
            "blocks": [],
            "metadata": {"slides": 0, "filename": filename, "error": "extraction_failed"},
        }

    return {
        "family": FileFamily.PRESENTATION,
        "title": title,
        "blocks": blocks,
        "metadata": {
            "slides": slide_count,
            "filename": filename,
            "extraction_mode": "xml-fallback",
        },
    }


def _shape_text_lines(shape: Any) -> list[str]:
    lines: list[str] = []

    if getattr(shape, "has_text_frame", False):
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            run_text = "".join((run.text or "") for run in paragraph.runs).strip()
            paragraph_text = run_text or (paragraph.text or "").strip()
            if not paragraph_text:
                continue
            split_lines = [line.strip() for line in paragraph_text.splitlines() if line.strip()]
            lines.extend(split_lines)

    if getattr(shape, "has_table", False):
        try:
            for row in shape.table.rows:
                cells = [str(cell.text or "").strip() for cell in row.cells if str(cell.text or "").strip()]
                if cells:
                    lines.append(" | ".join(cells))
        except Exception:
            pass

    if not lines and hasattr(shape, "text"):
        raw_text = str(getattr(shape, "text", "") or "").strip()
        if raw_text:
            lines.extend([line.strip() for line in raw_text.splitlines() if line.strip()])

    return _dedupe_preserve_order(lines)


def _extract_notes_text(slide: Any) -> str | None:
    try:
        if not slide.has_notes_slide:
            return None
        notes_frame = slide.notes_slide.notes_text_frame
        text = str(getattr(notes_frame, "text", "") or "").strip()
    except Exception:
        return None
    if not text or len(text) < 8:
        return None
    return _join_lines(text.splitlines())


def _read_xml_texts(*, zf: zipfile.ZipFile, xml_name: str) -> list[str]:
    try:
        root = ET.fromstring(zf.read(xml_name))
    except Exception:
        return []

    texts: list[str] = []
    for element in root.iter():
        if not element.tag.endswith("}t"):
            continue
        value = str(element.text or "").strip()
        if value:
            texts.append(value)
    return _dedupe_preserve_order(texts)


def _sorted_pptx_xml_names(names: Any) -> list[str]:
    def _index(name: str) -> int:
        match = re.search(r"(\d+)\.xml$", name)
        if not match:
            return 10**9
        return int(match.group(1))

    return sorted((str(name) for name in names), key=lambda name: (_index(name), name))


def _normalize_title(filename: str) -> str:
    normalized = str(filename or "").strip() or "apresentacao"
    normalized = normalized.rsplit("/", maxsplit=1)[-1].rsplit("\\", maxsplit=1)[-1]
    normalized = re.sub(r"\.pptx?$", "", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"[_-]+", " ", normalized).strip()
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized.title() if normalized else "Apresentacao"


def _join_lines(values: Any) -> str | None:
    lines = [str(value or "").strip() for value in values]
    compact = [line for line in lines if line]
    if not compact:
        return None
    return " ".join(compact)


def _dedupe_preserve_order(values: list[str]) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for raw in values:
        value = str(raw or "").strip()
        if not value:
            continue
        if value in seen:
            continue
        seen.add(value)
        output.append(value)
    return output
